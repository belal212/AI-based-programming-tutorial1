"""
Preprocessing Pipeline for Unstructured Data
This script extracts and normalizes data from various document types for LLM applications.
"""

import json
import os
from pathlib import Path

import pdfplumber
from docx import Document
import pandas as pd
from pptx import Presentation
from ebooklib import epub
import ebooklib


def extract_text_from_pdf(pdf_path):
    """Extract text from PDF files using pdfplumber."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
    return text.strip()


def extract_text_from_docx(docx_path):
    """Extract text from Word documents."""
    doc = Document(docx_path)
    text = " ".join([paragraph.text for paragraph in doc.paragraphs])
    return text.strip()


def extract_tables_from_excel(excel_path):
    """Extract tables from Excel files."""
    excel_data = pd.read_excel(excel_path, sheet_name=None)
    tables = {sheet: df.to_dict(orient='records') for sheet, df in excel_data.items()}
    return tables


def extract_text_from_ppt(ppt_path):
    """Extract text from PowerPoint presentations."""
    presentation = Presentation(ppt_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


def extract_text_from_epub(epub_path):
    """Extract text from EPUB files."""
    book = epub.read_epub(epub_path)
    text = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            text += item.get_content().decode('utf-8')
    return text.strip()


def normalize_to_json(file_type, content, metadata):
    """Normalize extracted data into JSON format."""
    return {
        "file_type": file_type,
        "content": content,
        "metadata": metadata
    }


def enrich_with_metadata(data, author=None, date=None, source=None):
    """Enrich data with metadata."""
    if author:
        data["metadata"]["author"] = author
    if date:
        data["metadata"]["date"] = date
    if source:
        data["metadata"]["source"] = source
    return data


def preprocess_document(file_path):
    """
    Main preprocessing pipeline for different document types.
    Extracts text and normalizes it into JSON format with metadata.
    """
    file_path = str(file_path)
    
    if file_path.endswith('.pdf'):
        text = extract_text_from_pdf(file_path)
        file_type = "PDF"
    elif file_path.endswith('.docx'):
        text = extract_text_from_docx(file_path)
        file_type = "Word"
    elif file_path.endswith('.xlsx'):
        tables = extract_tables_from_excel(file_path)
        text = str(tables)
        file_type = "Excel"
    elif file_path.endswith('.pptx'):
        text = extract_text_from_ppt(file_path)
        file_type = "PowerPoint"
    elif file_path.endswith('.epub'):
        text = extract_text_from_epub(file_path)
        file_type = "EPUB"
    else:
        raise ValueError(f"Unsupported file type: {file_path}")
    
    metadata = {
        "author": "Unknown",
        "date": "2023-10-01",
        "source": os.path.basename(file_path)
    }
    
    return normalize_to_json(file_type, text, metadata)


def process_folder(folder_path=".", output_folder="output"):
    """
    Process all supported documents in a folder and save as JSON files.
    """
    # Create output folder if it doesn't exist
    Path(output_folder).mkdir(exist_ok=True)
    
    # Supported file extensions
    extensions = ['.pdf', '.docx', '.xlsx', '.pptx', '.epub']
    
    # Process each file
    processed_files = []
    for file in Path(folder_path).iterdir():
        if file.is_file() and file.suffix in extensions:
            try:
                print(f"Processing {file.name}...")
                result = preprocess_document(file)
                
                # Save to output folder
                output_file = Path(output_folder) / f"{file.stem}_normalized.json"
                with open(output_file, 'w', encoding='utf-8') as f:
                    json.dump(result, f, indent=4, ensure_ascii=False)
                
                print(f"✓ Saved to {output_file}")
                processed_files.append(file.name)
            except Exception as e:
                print(f"✗ Error processing {file.name}: {e}")
    
    return processed_files


if __name__ == "__main__":
    print("=" * 50)
    print("Document Preprocessing Pipeline")
    print("=" * 50)
    
    # Process documents in the current directory
    processed = process_folder()
    
    if processed:
        print(f"\nSuccessfully processed {len(processed)} document(s):")
        for filename in processed:
            print(f"  - {filename}")
    else:
        print("\nNo supported documents found in the current directory.")
        print("Supported formats: PDF, DOCX, XLSX, PPTX, EPUB")
        print("\nExample usage:")
        print("  1. Place your documents in this directory")
        print("  2. Run: python preprocess_pipeline.py")
        print("  3. Check the 'output' folder for normalized JSON files")
